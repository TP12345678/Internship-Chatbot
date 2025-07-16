import os
import pandas as pd
from pptx import Presentation
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings

DATA_DIR = "./data"  

def extract_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    if "question" in df.columns and "answer" in df.columns:
        return (df["question"].astype(str) + " " + df["answer"].astype(str)).tolist()
    else:
        print(f"Skipping {file_path} — 'question' or 'answer' columns not found.")
        return []

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slide_texts = []
    for slide in prs.slides:
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        combined = "\n".join(text_parts).strip()
        if combined:
            slide_texts.append(combined)
    return slide_texts

# Automatically read all .pptx and .csv files
all_text_blocks = []
for filename in os.listdir(DATA_DIR):
    file_path = os.path.join(DATA_DIR, filename)
    if filename.lower().endswith(".csv"):
        print(f"📄 Reading CSV: {filename}")
        all_text_blocks.extend(extract_text_from_csv(file_path))
    elif filename.lower().endswith(".pptx"):
        print(f"📊 Reading PPTX: {filename}")
        all_text_blocks.extend(extract_text_from_pptx(file_path))
    else:
        print(f"⏭️ Skipping unsupported file: {filename}")

# Create embeddings and Chroma DB
embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")

vectordb = Chroma.from_texts(
    texts=all_text_blocks,
    embedding=embedding_model,
    persist_directory="chroma_db"
)

vectordb.persist() 
print("Vector DB created and persisted from all supported files in /data")
