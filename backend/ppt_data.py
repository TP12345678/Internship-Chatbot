from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from PIL import Image
import io

def extract_text_from_pptx_with_images(file_path):
    prs = Presentation(file_path)
    slide_texts = []

    for i, slide in enumerate(prs.slides):
        text_parts = []

        for shape in slide.shapes:
            # Text from text boxes
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

            # Text from images using OCR
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_stream = io.BytesIO(image_bytes)

                try:
                    img = Image.open(image_stream)
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        text_parts.append(ocr_text.strip())
                except Exception as e:
                    print(f"[Slide {i+1}] OCR failed on image: {e}")

        combined_text = "\n".join(text_parts).strip()
        if combined_text:
            slide_texts.append(combined_text)

    return slide_texts
