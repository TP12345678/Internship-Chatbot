import os
import pandas as pd
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from PIL import Image, UnidentifiedImageError
import io
import re 


def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error extracting text from PDF {pdf_path}: {e}")
    return text

def extract_text_from_image(image_path_or_bytes, source_info=""):
  
    text = ""
    img = None
    try:
        if isinstance(image_path_or_bytes, (str, os.PathLike)):
            img = Image.open(image_path_or_bytes)
        elif isinstance(image_path_or_bytes, bytes):
            img = Image.open(io.BytesIO(image_path_or_bytes))
        else:
            raise ValueError("Input must be a file path or bytes.")

        text = pytesseract.image_to_string(img)
    except pytesseract.TesseractNotFoundError:
        print("\nOCR Error")
        print("Tesseract OCR is not installed.")
        print("Please install Tesseract-OCR")
        print("-----------------\n")
    except UnidentifiedImageError:
        raise UnidentifiedImageError(
            f"Pillow cannot identify or load this image format for {source_info}. "
            "It might be a WMF, SVG, or a corrupted/unusual PNG. "
            "Please try converting it to a standard PNG/JPG."
        )
    except Exception as e:
        print(f"Error performing OCR on image {source_info}: {e}")
    finally:
        if img:
            img.close()
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide_idx, slide in enumerate(prs.slides):
            text += f"\n--- Slide {slide_idx + 1} ---\n"
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    source_info = f"slide {slide_idx + 1}, shape {shape_idx + 1} (bytes_len: {len(image_bytes)})"
                    try:
                        ocr_text = extract_text_from_image(image_bytes, source_info)
                        if ocr_text:
                            text += f"[OCR from image on {source_info.split(' (bytes_len:')[0]}]:\n{ocr_text}\n"
                    except UnidentifiedImageError as uie:
                        print(f"  Error processing image on {source_info}: {uie} Please convert it to PNG/JPG for OCR.")
                    except Exception as img_e:
                        print(f"  Error processing image on {source_info}: {img_e}")

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                if notes_text_frame and notes_text_frame.text:
                    text += f"\n--- Notes for Slide {slide_idx + 1} ---\n"
                    text += notes_text_frame.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX {pptx_path}: {e}")
    return text

def extract_text_from_csv(csv_path):
    text = ""
    try:
        df = pd.read_csv(csv_path)
        text = df.to_string(index=False)
    except Exception as e:
        print(f"Error extracting text from CSV {csv_path}: {e}")
    return text

def extract_text_from_txt(txt_path):
    text = ""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            text = f.read()
    except Exception as e:
        print(f"Error extracting text from TXT {txt_path}: {e}")
    return text

def split_document_by_case_study(full_text, source_path):
    case_study_splits = re.split(r'(\nCase Study : )', full_text)
    
    processed_splits = []
    current_segment = ""
    for i, segment in enumerate(case_study_splits):
        if segment.startswith('\nCase Study : '):
            if current_segment: 
                processed_splits.append(current_segment.strip())
            current_segment = segment 
        else:
            current_segment += segment
    
    if current_segment: 
        processed_splits.append(current_segment.strip())

    if len(processed_splits) <= 1 and not processed_splits[0].startswith('Case Study : '):
        return [{"text": full_text, "source": source_path}]

    documents_from_split = []
    for i, split_text in enumerate(processed_splits):
        if split_text: 
            match = re.search(r'Case Study : ([^\n]+)', split_text)
            case_study_name = match.group(1).strip() if match else f"Unnamed Case Study {i+1}"
            documents_from_split.append({
                "text": split_text,
                "source": f"{source_path} (Case Study: {case_study_name})"
            })
    return documents_from_split


def load_documents_from_folder(data_folder, single_file_path=None):

    documents = []
    if not os.path.exists(data_folder):
        print(f"Data folder '{data_folder}' not found. Please create it and add your data.")
        return documents

    image_extensions = [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"]

    files_to_process = []
    if single_file_path:
        files_to_process.append(single_file_path)
    else:
        for root, _, files in os.walk(data_folder):
            for file_name in files:
                files_to_process.append(os.path.join(root, file_name))

    for file_path in files_to_process:
        file_extension = os.path.splitext(file_path)[1].lower()
        extracted_text = ""

        if file_extension == ".pdf":
            extracted_text = extract_text_from_pdf(file_path)
        elif file_extension == ".pptx":
            extracted_text = extract_text_from_pptx(file_path)
        elif file_extension == ".csv":
            extracted_text = extract_text_from_csv(file_path)
        elif file_extension == ".txt":
            extracted_text = extract_text_from_txt(file_path)
        elif file_extension in image_extensions:
            source_info = f"file {file_path}"
            try:
                extracted_text = extract_text_from_image(file_path, source_info)
            except UnidentifiedImageError as uie:
                print(f"  Error processing image file {file_path}: {uie}")
            except Exception as e:
                print(f"  Error processing image file {file_path}: {e}")
        else:
            print(f"Skipping unsupported file type: {file_name}")
            continue

        if extracted_text:
            if "IDC Digital Corporate Deck- Middle East- V1 (1).txt" in file_path:
                print(f"Applying case study splitting for: {file_path}")
                split_docs = split_document_by_case_study(extracted_text, file_path)
                documents.extend(split_docs)
            else:
                documents.append({"text": extracted_text, "source": file_path})
            print(f"Extracted text from: {file_path}")
    return documents
