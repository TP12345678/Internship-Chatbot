import os
import pandas as pd
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from PIL import Image, UnidentifiedImageError
import io
import tempfile
import shutil

# --- Text Extraction Functions ---

def extract_text_from_pdf(pdf_path):
    """Extracts text from a PDF file."""
    text = ""
    try:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error extracting text from PDF {pdf_path}: {e}")
    return text

def extract_text_from_image(image_path_or_bytes, source_info=""):
    """
    Extracts text from an image file or bytes using OCR (Optical Character Recognition).
    Requires Tesseract-OCR to be installed on the system.
    """
    text = ""
    img = None
    try:
        if isinstance(image_path_or_bytes, (str, os.PathLike)):
            img = Image.open(image_path_or_bytes)
        elif isinstance(image_path_or_bytes, bytes):
            img = Image.open(io.BytesIO(image_path_or_bytes))
        else:
            raise ValueError("Input must be a file path or bytes.")

        # Perform OCR using pytesseract
        text = pytesseract.image_to_string(img)
    except pytesseract.TesseractNotFoundError:
        print("\n--- OCR Error ---")
        print("Tesseract-OCR is not installed or not in your system's PATH.")
        print("Please install Tesseract-OCR. See the main response for instructions.")
        print("-----------------\n")
        raise # Re-raise to ensure main process knows about the critical error
    except UnidentifiedImageError:
        # This specifically catches issues like unsupported image formats (e.g., WMF, SVG, or corrupted PNG)
        raise UnidentifiedImageError(
            f"Pillow cannot identify or load this image format for {source_info}. "
            "It might be a WMF, SVG, or a corrupted/unusual PNG. "
            "Please try converting it to a standard PNG/JPG."
        )
    except Exception as e:
        print(f"Error performing OCR on image {source_info}: {e}")
        # Optionally, re-raise if you want the main process to stop on any OCR error
        # raise
    finally:
        if img:
            img.close() # Close the image to free up resources
    return text

def extract_text_from_pptx(pptx_path):
    """
    Extracts text from a PPTX file (slides, notes, and performs OCR on embedded images).
    """
    text = ""
    print(f"  Processing PPTX: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_content = "" # Accumulate text for the current slide
            print(f"    --- Processing Slide {slide_idx + 1} ---")
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = f"Slide {slide_idx + 1}, Shape {shape_idx + 1} (Type: {shape.shape_type.name})"
                
                # Extract text from standard text shapes
                if hasattr(shape, "text"):
                    if shape.text.strip(): # Check if text is not just whitespace
                        slide_text_content += shape.text + "\n"
                        print(f"      Found text in {shape_info}: '{shape.text.strip()[:50]}...'")
                    else:
                        print(f"      Found empty text in {shape_info}.")
                # Check if the shape is a picture and perform OCR
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob # Get raw image data
                    temp_image_dir = None
                    try:
                        source_info = f"slide {slide_idx + 1}, shape {shape_idx + 1} (bytes_len: {len(image_bytes)})"
                        
                        ocr_text = extract_text_from_image(image_bytes, source_info)
                        if ocr_text.strip(): # Check if OCR actually found text
                            slide_text_content += f"[OCR from image on {source_info.split(' (bytes_len:')[0]}]:\n{ocr_text}\n"
                            print(f"      OCR extracted text from image in {shape_info}: '{ocr_text.strip()[:50]}...'")
                        else:
                            print(f"      OCR found no text in image in {shape_info}.")
                    except UnidentifiedImageError as uie:
                        # For unsupported image formats, try to save the raw blob for inspection
                        if image_bytes:
                            temp_image_dir = tempfile.mkdtemp(prefix="unsupported_img_")
                            temp_image_path = os.path.join(temp_image_dir, f"slide_{slide_idx+1}_shape_{shape_idx+1}_raw_image.bin")
                            with open(temp_image_path, 'wb') as f:
                                f.write(image_bytes)
                            print(f"  Error processing image on {shape_info}: {uie} Raw image data saved to: {temp_image_path}. Please inspect this file. Convert to PNG/JPG for OCR.")
                        else:
                            print(f"  Error processing image on {shape_info}: {uie} (No image bytes found).")
                    except Exception as img_e:
                        print(f"  Error processing image on {shape_info}: {img_e}")
                    finally:
                        # Clean up the temporary directory if it was created
                        if temp_image_dir and os.path.exists(temp_image_dir):
                            try:
                                shutil.rmtree(temp_image_dir)
                            except OSError as e:
                                print(f"Error deleting temporary directory {temp_image_dir}: {e}")
                else:
                    print(f"      Skipping non-text/non-picture shape: {shape_info}")

            # Extract text from notes (if any)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                if notes_text_frame and notes_text_frame.text.strip(): # Check if notes_text_frame and its text exist and is not empty
                    slide_text_content += f"\n--- Notes for Slide {slide_idx + 1} ---\n"
                    slide_text_content += notes_text_frame.text + "\n"
                    print(f"      Found notes text for Slide {slide_idx + 1}: '{notes_text_frame.text.strip()[:50]}...'")
                else:
                    print(f"      No notes text found for Slide {slide_idx + 1}.")
            
            # Only add accumulated slide text if it's not empty
            if slide_text_content.strip():
                text += slide_text_content + "\n"
            else:
                print(f"    Slide {slide_idx + 1} yielded no extractable text.")

    except Exception as e:
        print(f"Error extracting text from PPTX {pptx_path}: {e}")
    return text

def extract_text_from_csv(csv_path):
    """Extracts text from a CSV file by reading all cells."""
    text = ""
    try:
        df = pd.read_csv(csv_path)
        text = df.to_string(index=False) # Convert DataFrame to string
    except Exception as e:
        print(f"Error extracting text from CSV {csv_path}: {e}")
    return text

def extract_text_from_txt(txt_path):
    """Extracts text from a TXT file."""
    text = ""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            text = f.read()
    except Exception as e:
        print(f"Error extracting text from TXT {txt_path}: {e}")
    return text

def load_documents_from_folder(data_folder):
    """
    Loads and extracts text from all supported documents in the specified folder,
    including OCR for images.
    Returns a list of dictionaries, each containing 'text' and 'source'.
    """
    documents = []
    if not os.path.exists(data_folder):
        print(f"Data folder '{data_folder}' not found. Please create it and add your data.")
        return documents

    # Define supported image extensions for direct OCR
    image_extensions = [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"]

    for root, _, files in os.walk(data_folder):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            file_extension = os.path.splitext(file_name)[1].lower()
            extracted_text = ""

            print(f"\nProcessing file: {file_path}")

            if file_extension == ".pdf":
                extracted_text = extract_text_from_pdf(file_path)
            elif file_extension == ".pptx":
                extracted_text = extract_text_from_pptx(file_path)
            elif file_extension == ".csv":
                extracted_text = extract_text_from_csv(file_path)
            elif file_extension == ".txt":
                extracted_text = extract_text_from_txt(file_path)
            elif file_extension in image_extensions: # Handle standalone image files
                source_info = f"file {file_path}"
                try:
                    extracted_text = extract_text_from_image(file_path, source_info)
                except UnidentifiedImageError as uie:
                    print(f"  Error processing image file {file_path}: {uie}")
                except Exception as e:
                    print(f"  Error processing image file {file_path}: {e}")
            else:
                print(f"Skipping unsupported file type: {file_name}")
                continue

            if extracted_text.strip(): # Only add if actual text was extracted
                documents.append({"text": extracted_text, "source": file_path})
                print(f"Successfully extracted text from: {file_path}")
            else:
                print(f"No significant text extracted from: {file_path}")
    return documents
